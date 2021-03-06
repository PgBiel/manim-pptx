import os
import platform
import glob
import subprocess
from manimlib.constants import *
import manimlib.addon_helper
from manimlib.scene.scene_file_writer import *

try:
    import pptx
except ImportError(e):
    print("{0} addon failed to load:\n\tPlease install python-pptx to the {0} directory".format(Main.addon_info()['name']))
from pptx import Presentation


class Main(object):
    # Constants
    PACKAGE_NAME = "PPTX"
    CUR_DIR = os.path.join(ADDON_DIR, PACKAGE_NAME)
    TEMPLATE_DIR = os.path.join(CUR_DIR, "templates")
    TEMPLATE_PPTX = os.path.join(TEMPLATE_DIR, "template.pptx")
    EXAMPLE_PPTX = os.path.join(TEMPLATE_DIR, "powerpoint.pptx")
    TEMPORARY_DIR = os.path.join(CUR_DIR, "temporary")
    LOG_DIR = os.path.join(CUR_DIR, 'pptx.log')
    CONFIG = { }

    def __str__():
        return Main.PACKAGE_NAME

    def set_config(cfg):
        # Fires after the cli arguments are parsed in __init__.py for addons to have access to the current config
        Main.CONFIG = cfg

    def loaded():
        # Fires when the addon is first initialized by manim
        #Main.verbose("{0} addon loaded successfully".format(Main.addon_info()['name']))
        return True

    def on_rendered():
        # Fires when a video is finished rendering
        if Main.CONFIG["all_args"].save_to_pptx:
            Main.create_ppt()

    def create_ppt():
        SLD_BLANK = 6
        if not os.path.exists(Main.TEMPORARY_DIR):
            os.makedirs(Main.TEMPORARY_DIR)
        if os.path.exists(Main.LOG_DIR):
            os.remove(Main.LOG_DIR)
        prs = Presentation(Main.TEMPLATE_PPTX)
        
        # Figure out where the movie parts are saved
        PART_DIR = os.path.join(
            os.path.dirname(manimlib.addon_helper.movie_paths[0]), "partial_movie_files", Main.CONFIG['scene_names'][0]
        )
        Main.log_line("PART_DIR = " + PART_DIR)
        parts = glob.glob(os.path.join(PART_DIR, "*.mp4"))
        read_parts = PART_DIR

        # Go through each video part and copy it over to the temporary directory.
        # If anti-duplicate is on, every other part is combined with the one before it and then copied over to the temp directory
        if Main.CONFIG["all_args"].anti_dupli_pptx:
            Main.log_line("Anti-duplication manually enabled")
            read_parts = Main.TEMPORARY_DIR
            for file in parts:
                i = int(Main.get_name(file))
                if not i % 2 == 0:
                    Main.log_line("Merging parts {} and {}...".format(str(i), str(i-1)))
                    merged_clip = Main.merge_videos(parts[i-1], file, os.path.join(Main.TEMPORARY_DIR, str(i-1).zfill(5) + ".mp4"))
                    Main.log_line("Merged to " + merged_clip)

        save_dir = os.path.join(os.path.dirname(manimlib.addon_helper.movie_paths[0]), Main.CONFIG['scene_names'][0] + ".pptx")
        slide_layout = prs.slide_layouts[SLD_BLANK]
        for file in glob.glob(os.path.join(read_parts, "*.mp4")):
            Main.log_line("Using file at " + file)
            # Load the example presentation and its timing element
            prs_ex = Presentation(Main.EXAMPLE_PPTX)
            timing_ex = prs_ex.slides[0].element[2]
            Main.log_line("\tGrabbed timing element, timing_ex = " + str(timing_ex))
            # Create a new slide
            slide = prs.slides.add_slide(slide_layout)
            # Generate video thumbnail
            thumb_file = os.path.join(Main.TEMPORARY_DIR, Main.get_name(file) + ".png")
            Main.log_line("\tGenerating video thumbnail...")
            Main.get_video_thumb(file, thumb_file)
            Main.log_line("\tThumbnail saved at " + thumb_file)

            # Add the video to the slide
            clip = slide.shapes.add_movie(file, 0, 0, prs.slide_width, prs.slide_height, mime_type='video/mp4', poster_frame_image=thumb_file)
            Main.log_line("\tAdded clip to slide")

            # Play the clip in fullscreen when the slide starts
            ## Get the id of the movie object we just added
            id = clip.element[0][0].attrib.get("id")
            Main.log_line("\tClip id = " + id)
            ## Make a copy of the timing element from the manually created PPTX,
            ## then change every spid to the clip id
            Main.log_line("\tUsing timing_ex as template...")
            timing = timing_ex
            timing[0][0][0][0][0][0][0][0][0][1][0][0][1][0][0][1][0][0][1][0].attrib["spid"] = id
            timing[0][0][0][0][1][0][1][0].attrib["spid"] = id
            timing[0][0][0][0][2][0][0][0][0][0].attrib["spid"] = id
            timing[0][0][0][0][2][0][2][0][0][1][0][0][1][0][0][1][0][0][1][0].attrib["spid"] = id
            timing[0][0][0][0][2][1][0][0][0].attrib["spid"] = id
            slide.element[2] = timing
            Main.log_line("\tAdded timing to slide, timing = " + str(timing))
            prs.save(save_dir)
            Main.log_line("\tPPTX saved to " + save_dir)
        
        Main.log_line("Final presentation saved to " + save_dir)
        print("\nPresentation ready at " + save_dir)
        if Main.CONFIG["all_args"].preview:
            Main.open_file(save_dir)

    def get_name(filename):
        pre, ext = os.path.splitext(filename)
        return pre.split(os.sep)[-1]

    def get_video_thumb(filename, imgname):
        command = [
                FFMPEG_BIN,
                '-y',  # overwrite output file if it exists
                '-loglevel', 'error',
                '-i', filename,
                '-vframes', '1',  # take only one frame
                imgname,
            ]
        subprocess.run(command, stdout=subprocess.PIPE)
        return imgname

    def get_frame_count(filename):
        command = [
            FFPROBE_BIN,
            '-v', 'error',
            '-count_frames',
            '-select_streams', 'v:0',
            '-show_entries', 'stream=nb_read_frames',
            '-of', 'default=nokey=1:noprint_wrappers=1',
            filename
        ]
        result = subprocess.run(command, stdout=subprocess.PIPE).stdout.decode('utf-8')
        return int(result)

    # TODO: Make this work
    def get_middle_video_frame(filename, imgname):
        command = [
            FFMPEG_BIN,
            '-y',  # overwrite output file if it exists
            '-loglevel', 'error',
            '-i', filename,
            '-vframes', '1',  # take only one frame
            imgname,
        ]
        subprocess.run(command, stdout=subprocess.PIPE)
        return imgname

    def merge_videos(clip1, clip2, output):
        vid_list = os.path.join(Main.TEMPORARY_DIR, "cliplist.txt") 
        with open(vid_list, 'w') as file:
            file.write("file '{}'\n".format(clip1))
            file.write("file '{}'".format(clip2))
        commands = [
            FFMPEG_BIN,
            '-y',  # overwrite output file if it exists
            '-f', 'concat',
            '-safe', '0',
            '-i', vid_list,
            '-loglevel', 'error',
            '-c', 'copy', output
        ]
        subprocess.run(commands, stdout=subprocess.PIPE)
        return output

    def open_file(file_path):
        # Taken from open_file_if_needed()
        current_os = platform.system()
        if current_os == "Windows":
            os.startfile(file_path)
        else:
            commands = []
            if current_os == "Linux":
                commands.append("xdg-open")
            elif current_os.startswith("CYGWIN"):
                commands.append("cygstart")
            else:  # Assume macOS
                commands.append("open")

            if config["show_file_in_finder"]:
                commands.append("-R")

            commands.append(file_path)

            # commands.append("-g")
            FNULL = open(os.devnull, 'w')
            sp.call(commands, stdout=FNULL, stderr=sp.STDOUT)
            FNULL.close()

    def log_line(text):
        Main.log_text(text.__str__() + "\n")

    def log_text(text):
        with open(Main.LOG_DIR, 'a') as the_file:
            the_file.write(text.__str__())

    def parser_args():
        # Return any command line flags that the addon adds
        return [
            {
                'flag': "--save_to_pptx",
                'action': "store_true",
                'help': "[{0}] Render the animations to a PowerPoint presentation".format(Main.PACKAGE_NAME)
            },
            {
                'flag': "--anti_dupli_pptx",
                'action': "store_true",
                'help': "[{0}] When exporting to PowerPoint, only use every other movie part".format(Main.PACKAGE_NAME)
            }
        ]

    def addon_info():
        return {
            'author': "Joshua \"Yoshi\" Askharoun",
            'name': Main.PACKAGE_NAME,
            'version' : "1.0.1.0",
            'desc': "Adds --save_to_pptx, which generates a PowerPoint with a slide for each animation"
        }